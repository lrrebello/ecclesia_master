import os
import PyPDF2
from docx import Document
from pptx import Presentation
import requests
from io import BytesIO

def extract_text_from_pdf(file_content):
    pdf_file = BytesIO(file_content)
    reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(file_content):
    docx_file = BytesIO(file_content)
    doc = Document(docx_file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_pptx(file_content):
    pptx_file = BytesIO(file_content)
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(file_path_or_url):
    try:
        if file_path_or_url.startswith(('http://', 'https://')):
            response = requests.get(file_path_or_url)
            content = response.content
            filename = file_path_or_url.split('/')[-1]
        else:
            with open(file_path_or_url, 'rb') as f:
                content = f.read()
            filename = os.path.basename(file_path_or_url)

        ext = filename.split('.')[-1].lower()
        
        if ext == 'pdf':
            return extract_text_from_pdf(content)
        elif ext == 'docx':
            return extract_text_from_docx(content)
        elif ext == 'pptx':
            return extract_text_from_pptx(content)
        elif ext in ['txt', 'md']:
            return content.decode('utf-8')
        else:
            return ""
    except Exception as e:
        print(f"Erro na extração: {e}")
        return ""
